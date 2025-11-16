from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# --- CONSTANTES DE DISEÑO ---
COURSE_FOOTER = (
    "CasoPractico2 - AliciaBeamud -  ‘Programación con Inteligencia Artificial: Introducción y gestión con LLMs –"
)
TITLE_COLOR = RGBColor(20, 40, 80)      # Azul oscuro
ACCENT = RGBColor(71, 134, 226)         # Azul didáctico
ACCENT_SOFT = RGBColor(227, 238, 255)   # Azul muy claro
TEXT_COLOR = RGBColor(40, 40, 40)
MUTED = RGBColor(110, 120, 130)


# --- FUNCIONES DE AYUDA ---

def add_footer(prs, slide, page_num=None):
    """Añade un pie de página institucional y número de diapositiva."""
    width, height = prs.slide_width, prs.slide_height
    tx = slide.shapes.add_textbox(Inches(0.5), height - Inches(0.5), width - Inches(1.0), Inches(0.3))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = COURSE_FOOTER
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT

    # Número de diapositiva a la derecha
    if page_num is not None:
        tx2 = slide.shapes.add_textbox(width - Inches(1.1), height - Inches(0.5), Inches(0.6), Inches(0.3))
        tf2 = tx2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = str(page_num)
        p2.font.size = Pt(10)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.RIGHT


# --- DIAPOSITIVA 1: PORTADA (FINAL CORREGIDA) ---
def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # en blanco
    w, h = prs.slide_width, prs.slide_height

    # Banda superior azul (fondo): AUMENTADA A 2.0 PULGADAS
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, w, Inches(2.0) 
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    # Título: POSICIÓN Y AJUSTADA
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), w - Inches(1.4), Inches(1.1)) 
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtítulo / curso: POSICIÓN Y AJUSTADA
    tx2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), w - Inches(1.4), Inches(0.6))
    tf2 = tx2.text_frame
    tf2.clear()
    s = tf2.paragraphs[0]
    s.text = subtitle
    s.font.size = Pt(16)
    s.font.color.rgb = RGBColor(245, 245, 245)

    add_footer(prs, slide, page_num=1)


# --- DIAPOSITIVAS DE CONTENIDO (CON WORD_WRAP) ---
def bullet_slide(prs, title, bullets, is_index=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    # Encabezado
    header = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), w - Inches(1.4), Inches(0.8))
    tf = header.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Línea de acento
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.2), Inches(3.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Caja de contenido (Fondo blanco)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), w - Inches(1.4), h - Inches(2.4)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = ACCENT

    # Viñetas: Word_wrap aplicado
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), w - Inches(2.6), h - Inches(2.8))
    tf2 = tx.text_frame
    tf2.clear()
    tf2.word_wrap = True 
    
    font_size = Pt(20) if is_index else Pt(20)
    
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = font_size
        p.font.color.rgb = TEXT_COLOR

    add_footer(prs, slide, page_num=len(prs.slides))


def section_break_slide(prs, title):
    """Añade una diapositiva de impacto visual como separador de sección."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    # Fondo de color ACCENT
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
    background.fill.solid()
    background.fill.fore_color.rgb = ACCENT
    background.line.fill.background()

    # Título centrado
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), w - Inches(1.0), Inches(2.0))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255) # Texto blanco
    p.alignment = PP_ALIGN.CENTER
    
    # Nota del curso abajo
    tx2 = slide.shapes.add_textbox(Inches(0.5), h - Inches(1.5), w - Inches(1.0), Inches(0.5))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = COURSE_FOOTER
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(230, 230, 230)
    p2.alignment = PP_ALIGN.CENTER


def two_columns_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    header = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), w - Inches(1.4), Inches(0.8))
    p = header.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = TITLE_COLOR

    # Línea de acento
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.2), Inches(3.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Columnas
    col_w = (w - Inches(2.2)) / 2
    # Izquierda
    box_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), col_w, h - Inches(2.4))
    box_l.fill.solid(); box_l.fill.fore_color.rgb = RGBColor(255, 255, 255); box_l.line.color.rgb = ACCENT

    txl = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), col_w - Inches(0.4), Inches(0.5))
    p = txl.text_frame.paragraphs[0]; p.text = left_title; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = TITLE_COLOR

    list_l = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), col_w - Inches(0.7), h - Inches(3.0))
    tf = list_l.text_frame; tf.clear()
    tf.word_wrap = True 
    for i, it in enumerate(left_items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = it; para.font.size = Pt(18); para.font.color.rgb = TEXT_COLOR

    # Derecha
    x_r = Inches(1.3) + col_w
    box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_r, Inches(1.5), col_w, h - Inches(2.4))
    box_r.fill.solid(); box_r.fill.fore_color.rgb = RGBColor(255, 255, 255); box_r.line.color.rgb = ACCENT

    txr = slide.shapes.add_textbox(x_r + Inches(0.2), Inches(1.7), col_w - Inches(0.4), Inches(0.5))
    p = txr.text_frame.paragraphs[0]; p.text = right_title; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = TITLE_COLOR

    list_r = slide.shapes.add_textbox(x_r + Inches(0.2), Inches(2.2), col_w - Inches(0.7), h - Inches(3.0))
    tf2 = list_r.text_frame; tf2.clear()
    tf2.word_wrap = True 
    for i, it in enumerate(right_items):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = it; para.font.size = Pt(18); para.font.color.rgb = TEXT_COLOR

    add_footer(prs, slide, page_num=len(prs.slides))


# --- FUNCIÓN PRINCIPAL DE CONSTRUCCIÓN ---
def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1) Portada
    title_slide(
        prs,
        title="El Reglamento Europeo de IA (AI Act)",
        subtitle="Un Resumen de la Nueva Legislación y su Impacto Regulatorio",
    )

    # 2) Índice
    bullet_slide(
        prs,
        "Índice de Contenidos",
        [
            "1. Origen y Objetivos del AI Act",
            "2. El Enfoque Basado en Riesgos",
            "3. Sistemas de Riesgo Inaceptable (Prohibidos)",
            "4. Sistemas de Alto Riesgo (Obligaciones)",
            "5. Regulación de Modelos de IA Generativa",
            "6. Implementación y Sanciones",
            "7. Conclusiones y Futuro",
        ],
        is_index=True
    )

    # 3) Origen y Objetivos
    bullet_slide(
        prs,
        "1. Origen y Objetivos del AI Act",
        [
            "La primera ley integral de IA del mundo, con alcance extraterritorial ('Efecto Bruselas').",
            "Objetivo principal: Fomentar la adopción de IA centrada en el ser humano, garantizando la seguridad y los derechos fundamentales.",
            "Asegurar que la IA en la UE sea segura, legal, ética y fiable.",
            "Busca establecer reglas uniformes para operadores y desarrolladores de IA.",
        ],
    )
    
    # 4) Enfoque Basado en Riesgos
    bullet_slide(
        prs,
        "2. El Enfoque Basado en Riesgos",
        [
            "Clasificación jerárquica de los sistemas de IA en 4 niveles según su capacidad de causar daño.",
            "Cuanto mayor es el riesgo, más estrictas son las obligaciones de cumplimiento.",
            "Niveles: Inaceptable (Prohibido), Alto, Limitado (Transparencia) y Mínimo (No regulado).",
            "La mayoría de las IAs (riesgo mínimo) no tendrán obligaciones adicionales.",
        ],
    )

    # DIAPOSITIVA DE SECCIÓN PARA IMPACTO VISUAL
    section_break_slide(prs, "El Núcleo del AI Act: Sistemas de Alto Riesgo")

    # 5) Riesgo Inaceptable (Prohibido)
    bullet_slide(
        prs,
        "3. Sistemas de Riesgo Inaceptable (Prohibidos)",
        [
            "Sistemas que representan una clara amenaza a los derechos fundamentales y la democracia.",
            "Prohibiciones clave:",
            "  •  Sistemas de puntuación social ('Social Scoring') por autoridades públicas.",
            "  •  Técnicas de manipulación subliminal o explotadora para causar daño.",
            "  •  Vigilancia biométrica masiva en tiempo real (con excepciones limitadas).",
        ],
    )

    # 6) Sistemas de Alto Riesgo (Obligaciones)
    two_columns_slide(
        prs,
        title="4. Sistemas de Alto Riesgo: Requisitos",
        left_title="Sectores Típicos",
        left_items=[
            "Dispositivos médicos y seguridad de productos.",
            "Infraestructuras críticas (agua, gas, electricidad).",
            "Educación y Recursos Humanos.",
            "Aplicación de la ley (policía y justicia).",
            "Gestión de migración y asilo.",
        ],
        right_title="Obligaciones del Proveedor",
        right_items=[
            "Implementar Sistemas de Gestión de Riesgos.",
            "Alta Calidad de los Datos de entrenamiento.",
            "Registro de Eventos (Logs) y Trazabilidad.",
            "Supervisión Humana Obligatoria.",
            "Evaluación de Conformidad (CE Mark) antes del lanzamiento.",
        ],
    )

    # 7) IA Generativa (Modelos Fundacionales)
    bullet_slide(
        prs,
        "5. Regulación de Modelos de IA Generativa",
        [
            "Los Modelos Fundacionales (ej. GPT, LaMDA) tienen nuevas obligaciones de transparencia.",
            "Obligaciones clave para desarrolladores:",
            "  •  Resumir y publicar los datos utilizados para el entrenamiento (respetando derechos de autor).",
            "  •  Diseñar el modelo para evitar la generación de contenido ilegal.",
            "  •  Etiquetar contenido generado por IA (Deepfakes) como artificial.",
        ],
    )

    # 8) Implementación y Sanciones
    bullet_slide(
        prs,
        "6. Implementación y Sanciones",
        [
            "Aplicación gradual a partir de 2024 (prohibiciones) y plena aplicación esperada para 2026-2027.",
            "Se crea la Oficina Europea de IA para supervisar la aplicación y guiar a los desarrolladores.",
            "Sanciones por incumplimiento:",
            "  •  Hasta 35 millones de euros o 7% de la facturación global (por violar sistemas prohibidos).",
            "  •  Hasta 15 millones de euros o 3% (por violar reglas de Alto Riesgo).",
        ],
    )

    # 9) Cierre y Conclusiones
    bullet_slide(
        prs,
        "Conclusiones: Impacto y Futuro",
        [
            "El AI Act está redefiniendo el panorama legal de la IA a nivel mundial.",
            "La clasificación por riesgo exige un inventario y evaluación proactiva de sistemas por parte de las organizaciones.",
            "La clave es la **trazabilidad** y la **transparencia**.",
            "Este marco busca fomentar la confianza en la tecnología sin sofocar la innovación responsable.",
        ],
    )

    file_name = "CasoPractico2_AliciaB_Reglamento_IA_Act.pptx"
    prs.save(file_name)
    print(f"✔ Presentación generada: {file_name}")


if __name__ == "__main__":
    build_presentation()
